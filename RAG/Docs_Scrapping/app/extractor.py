import os
import json
import pandas as pd
from typing import Tuple
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import pytesseract
from pdf2image import convert_from_path
from PIL import Image


SUPPORTED_TEXT_TYPES = {".txt", ".md", ".log", ".py", ".json", ".csv", ".rtf"}


def extract_text_from_file(path: str) -> Tuple[str, str]:
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(path)

    elif ext == ".docx":
        return "docx", _extract_docx(path)

    elif ext == ".pptx":
        return "pptx", _extract_pptx(path)

    elif ext in [".xlsx", ".xls"]:
        return "excel", _extract_excel(path)

    elif ext == ".csv":
        return "csv", _extract_csv(path)

    elif ext in [".html", ".htm"]:
        return "html", _extract_html(path)

    elif ext in SUPPORTED_TEXT_TYPES:
        return "text", _extract_text(path)

    elif ext in [".png", ".jpg", ".jpeg"]:
        return "image", _extract_image_ocr(path)

    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ---------------- PDF ---------------- #

def _extract_pdf(path: str) -> str:
    reader = PdfReader(path)
    text = "\n".join(page.extract_text() or "" for page in reader.pages)

    if text.strip():
        return text

    # OCR fallback
    images = convert_from_path(path)
    ocr_text = "\n".join(
        pytesseract.image_to_string(img) for img in images
    )
    return ocr_text


# ---------------- DOCX ---------------- #

def _extract_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


# ---------------- PPTX ---------------- #

def _extract_pptx(path: str) -> str:
    prs = Presentation(path)
    slides_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)

    return "\n".join(slides_text)


# ---------------- EXCEL ---------------- #

def _extract_excel(path: str) -> str:
    xls = pd.ExcelFile(path)
    sheets_text = []

    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        sheets_text.append(df.to_string(index=False))

    return "\n".join(sheets_text)


# ---------------- CSV ---------------- #

def _extract_csv(path: str) -> str:
    df = pd.read_csv(path)
    return df.to_string(index=False)


# ---------------- HTML ---------------- #

def _extract_html(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f, "lxml")
        return soup.get_text(separator="\n")


# ---------------- TEXT ---------------- #

def _extract_text(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="latin-1") as f:
            return f.read()


# ---------------- IMAGE OCR ---------------- #

def _extract_image_ocr(path: str) -> str:
    img = Image.open(path)
    return pytesseract.image_to_string(img)
