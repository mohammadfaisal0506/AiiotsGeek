import os, io, tempfile
from typing import Tuple, List
from pdfminer.high_level import extract_text as pdfminer_extract
import fitz  
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from .utils import logger

def extract_text_from_pdf(file_path: str) -> str:
   
    try:
        doc = fitz.open(file_path)
        pages = []
        for page in doc:
            text = page.get_text("text")
            if text and text.strip():
                pages.append(text)
            else:
                # attempt OCR on page image if no text
                pix = page.get_pixmap()
                img = Image.open(io.BytesIO(pix.tobytes()))
                ocr = pytesseract.image_to_string(img)
                pages.append(ocr)
        return "\n".join(pages)
    except Exception as e:
        logger.warning(f"PyMuPDF failed: {e}. Falling back to pdfminer.")
    # fallback
    try:
        return pdfminer_extract(file_path)
    except Exception as e:
        logger.exception("pdfminer fallback failed.")
        raise

def extract_text_from_docx(file_pathpath: str) -> str:
    doc = Document(file_pathpath)
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(paragraphs)

def extract_text_from_pptx(file_pathpath: str) -> str:
    pres = Presentation(file_pathpath)
    texts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text
                if t and t.strip():
                    texts.append(t)
    return "\n".join(texts)

def extract_text_from_txt(file_pathpath: str) -> str:
    with open(file_pathpath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_file(file_pathpath: str) -> Tuple[str, str]:
    ext = os.file_pathpath.splitext(file_pathpath)[1].lower()
    if ext == ".pdf":
        return "pdf", extract_text_from_pdf(file_pathpath)
    if ext == ".docx":
        return "docx", extract_text_from_docx(file_pathpath)
    if ext in [".pptx", ".ppt"]:
        return "pptx", extract_text_from_pptx(file_pathpath)
    if ext in [".txt", ".md"]:
        return "txt", extract_text_from_txt(file_pathpath)
    # try to open as image for OCR
    try:
        img = Image.open(file_pathpath)
        txt = pytesseract.image_to_string(img)
        return "image", txt
    except Exception:
        raise ValueError(f"Unsupported file type: {ext}")
